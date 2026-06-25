import streamlit as st
from openai import OpenAI
import docx
from docx.shared import Pt, RGBColor, Inches
from docx.oxml import parse_xml  
from docx.oxml.ns import nsdecls, qn  
from pptx import Presentation
from io import BytesIO
import pypdf
import fitz  # PyMuPDF，用于将PDF转图片
import base64

# 1. 初始化客户端
# 从 Streamlit 环境变量中安全读取 API Key
DEEPSEEK_KEY = st.secrets["DEEPSEEK_API_KEY"]
QWEN_KEY = st.secrets["QWEN_API_KEY"]

client_deepseek = OpenAI(
    api_key=DEEPSEEK_KEY, 
    base_url="https://api.deepseek.com/v1"
)

client_vision = OpenAI(
    api_key=QWEN_KEY,
    base_url="https://dashscope.aliyuncs.com/compatible-mode/v1"
)

# 初始化 Session State
if 'generated_plan' not in st.session_state:
    st.session_state['generated_plan'] = ""
if 'generated_outline' not in st.session_state:
    st.session_state['generated_outline'] = ""
if 'generated_quiz' not in st.session_state:
    st.session_state['generated_quiz'] = ""

def encode_image_to_base64(image_bytes):
    return base64.b64encode(image_bytes).decode('utf-8')

# 1.2 多模态视觉解析函数
def parse_image_with_vision(image_bytes, description="这是一张课本页面"):
    try:
        base64_image = encode_image_to_base64(image_bytes)
        response = client_vision.chat.completions.create(
            model="qwen-vl-plus",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": f"{description}。请：1. 提取所有文字。2. 描述插图、场景和排版。3. 总结核心教学点。"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=1500
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"视觉解析失败: {e}"

# 2. 增强版解析引擎（支持扫描版PDF）
@st.cache_data
def extract_text_cached(file_name, file_bytes, file_type, range_type, start, end, use_vision_for_pdf=False):
    extracted_text = ""
    start_idx = max(0, start - 1)
    file_stream = BytesIO(file_bytes)
    
    try:
        if file_type in ['png', 'jpg', 'jpeg']:
            extracted_text = parse_image_with_vision(file_bytes, "这是一张图片课件")
            
        elif file_type == 'pdf':
            if use_vision_for_pdf:
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                total = len(doc)
                end_idx = min(total, end) if range_type == "自定义范围" else total
                
                results = []
                for i in range(start_idx, end_idx):
                    page = doc.load_page(i)
                    pix = page.get_pixmap(matrix=fitz.Matrix(1.3, 1.3)) 
                    img_bytes = pix.tobytes("jpg")
                    res = parse_image_with_vision(img_bytes, f"这是PDF的第{i+1}页扫描件")
                    results.append(f"\n--- PDF第{i+1}页视觉解析 ---\n{res}")
                extracted_text = "\n".join(results)
                doc.close()
            else:
                pdf_reader = pypdf.PdfReader(file_stream)
                total = len(pdf_reader.pages)
                end_idx = min(total, end) if range_type == "自定义范围" else total
                for i in range(start_idx, end_idx):
                    extracted_text += f"\n--- Page {i+1} ---\n" + (pdf_reader.pages[i].extract_text() or "")
        
        elif file_type == 'pptx':
            prs = Presentation(file_stream)
            total = len(prs.slides)
            end_idx = min(total, end) if range_type == "自定义范围" else total
            ppt_text_list = []
            for i in range(start_idx, end_idx):
                slide = prs.slides[i]
                ppt_text_list.append(f"\n--- Slide {i + 1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        ppt_text_list.append(shape.text.strip())
            extracted_text = "\n".join(ppt_text_list)
            
        elif file_type == 'docx':
            doc = docx.Document(file_stream)
            total = len(doc.paragraphs)
            end_idx = min(total, end) if range_type == "自定义范围" else total
            extracted_text = "\n".join([para.text for para in doc.paragraphs[start_idx:end_idx]])
            
    except Exception as e:
        return f"ERROR: 解析出错 - {str(e)}"
        
    return extracted_text.strip()

# 3. 将 Markdown 转换为 Word 的辅助函数
def convert_markdown_to_docx(lesson_plan, outline, quiz, title="教学设计方案"):
    doc = docx.Document()
    
    # 页面边距设置（标准 A4 公文边距）
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)
        
    font_family_name = 'Microsoft YaHei'  
    
    # 强行设置东亚（中文）和西文属性，解决 Word 回退字体的 Bug
    def set_run_font_and_format(run, font_name=font_family_name, size_pt=10.5, color_rgb=None, bold=None):
        run.font.name = font_name
        rPr = run._r.get_or_add_rPr()
        rFonts = rPr.get_or_add_rFonts()
        rFonts.set(qn('w:eastAsia'), font_name)
        
        if size_pt:
            run.font.size = Pt(size_pt)
        if color_rgb:
            run.font.color.rgb = color_rgb
        if bold is not None:
            run.bold = bold

    # 用 XML 修改单元格底色（实现高级表格视觉）
    def set_cell_background(cell, fill_hex):
        shading_elm = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{fill_hex}"/>')
        cell._tc.get_or_add_tcPr().append(shading_elm)

    def add_styled_paragraph(text, size_pt=10.5, color_rgb=RGBColor(0x33, 0x33, 0x33), bold=False, heading_level=None):
        if heading_level is not None:
            p = doc.add_heading(level=heading_level)
        else:
            p = doc.add_paragraph()
            
        p.paragraph_format.line_spacing = 1.25
        p.paragraph_format.space_after = Pt(6)
        
        run = p.add_run(text)
        set_run_font_and_format(run, font_family_name, size_pt, color_rgb, bold)
        return p

    # Markdown 实体表格渲染引擎
    def render_markdown_table_to_docx(table_lines):
        clean_rows = []
        for line in table_lines:
            cells = [c.strip() for c in line.split('|')[1:-1]]
            if not cells:
                continue
            is_separator = all(all(char in '-: ' for char in cell) for cell in cells)
            if not is_separator:
                clean_rows.append(cells)
                
        if not clean_rows:
            return
            
        num_rows = len(clean_rows)
        num_cols = len(clean_rows[0])
        
        table = doc.add_table(rows=num_rows, cols=num_cols)
        table.style = 'Table Grid'
        
        for r_idx, row_data in enumerate(clean_rows):
            row = table.rows[r_idx]
            is_header = (r_idx == 0)
            
            for c_idx, cell_value in enumerate(row_data):
                if c_idx < len(row.cells):
                    cell = row.cells[c_idx]
                    cell.text = "" 
                    p = cell.paragraphs[0]
                    p.paragraph_format.line_spacing = 1.15
                    p.paragraph_format.space_before = Pt(4)
                    p.paragraph_format.space_after = Pt(4)
                    
                    run = p.add_run(cell_value)
                    
                    if is_header:
                        set_run_font_and_format(run, font_family_name, size_pt=10, color_rgb=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
                        set_cell_background(cell, "1B365D")
                    else:
                        set_run_font_and_format(run, font_family_name, size_pt=9.5, color_rgb=RGBColor(0x33, 0x33, 0x33), bold=False)
                        if r_idx % 2 == 0:
                            set_cell_background(cell, "F2F4F7")

    # 模块段落与表格混合写入引擎
    def write_styled_section(main_title, content_text):
        add_styled_paragraph(main_title, size_pt=16, color_rgb=RGBColor(0x1B, 0x36, 0x5D), bold=True)
        
        lines = content_text.split('\n')
        i = 0
        n = len(lines)
        
        while i < n:
            line = lines[i]
            stripped_line = line.strip()
            
            if stripped_line.startswith('|') and stripped_line.endswith('|'):
                table_lines = []
                while i < n and lines[i].strip().startswith('|') and lines[i].strip().endswith('|'):
                    table_lines.append(lines[i].strip())
                    i += 1
                render_markdown_table_to_docx(table_lines)
                continue 
                
            if not stripped_line:
                i += 1
                continue
            
            if stripped_line.startswith('# '):
                add_styled_paragraph(stripped_line[2:], heading_level=1, size_pt=14, color_rgb=RGBColor(0x1B, 0x36, 0x5D), bold=True)
            elif stripped_line.startswith('## '):
                add_styled_paragraph(stripped_line[3:], heading_level=2, size_pt=12.5, color_rgb=RGBColor(0x4B, 0x6B, 0x94), bold=True)
            elif stripped_line.startswith('### '):
                add_styled_paragraph(stripped_line[4:], heading_level=3, size_pt=11, color_rgb=RGBColor(0x33, 0x33, 0x33), bold=True)
            elif stripped_line.startswith('- ') or stripped_line.startswith('* '):
                p = doc.add_paragraph(style='List Bullet')
                p.paragraph_format.line_spacing = 1.25
                p.paragraph_format.space_after = Pt(4)
                run = p.add_run(stripped_line[2:])
                set_run_font_and_format(run, font_family_name, size_pt=10.5, color_rgb=RGBColor(0x33, 0x33, 0x33), bold=False)
            elif stripped_line.startswith('1. ') or stripped_line.startswith('2. ') or stripped_line.startswith('3. '):
                p = doc.add_paragraph(style='List Number')
                p.paragraph_format.line_spacing = 1.25
                p.paragraph_format.space_after = Pt(4)
                run = p.add_run(stripped_line)
                set_run_font_and_format(run, font_family_name, size_pt=10.5, color_rgb=RGBColor(0x33, 0x33, 0x33), bold=False)
            else:
                add_styled_paragraph(stripped_line, size_pt=10.5, color_rgb=RGBColor(0x33, 0x33, 0x33), bold=False)
            
            i += 1
                
    write_styled_section("第一部分：详细教案 (Lesson Plan)", lesson_plan)
    doc.add_page_break()
    write_styled_section("第二部分：配套授课大纲 (Lecture Outline)", outline)
    doc.add_page_break()
    write_styled_section("第三部分：随堂测试与互动 (Quiz & Activities)", quiz)
            
    bio = BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio

# 4. 页面基本配置
st.set_page_config(
    page_title="全学段英语智能教案生成器 (多模态版)",
    page_icon="📚",
    layout="wide"
)

st.title("📚 全学段英语智能教案生成器")
st.write("解析 PDF/Word/PPT 以及**图片课件与插画**，自动根据课堂总时长调整教学深度与配套大纲。")

# 5. 左右布局
col1, col2 = st.columns([1, 1.3])

with col1:
    st.subheader("🛠️ 教案与大纲参数配置")
    
    course_name = st.text_input("1. 课程名称", value="英语")
    
    school_stage = st.selectbox(
        "2. 选择授课学段",
        ["小学", "初中", "高中", "大学"]
    )
    
    if school_stage == "小学":
        grade_level = st.selectbox(
            "具体年级段", 
            ["低段（一、二年级 - 趣味听说）", "中段（三、四年级 - 读写萌芽）", "高段（五、六年级 - 基础读写技能）"]
        )
    elif school_stage == "初中":
        grade_level = st.selectbox(
            "具体年级", 
            ["初一（七年级）", "初二（八年级）", "初三（九年级）"]
        )
    elif school_stage == "高中":
        grade_level = st.selectbox(
            "具体年级", 
            ["高一", "高二", "高三"]
        )
    elif school_stage == "大学":
        grade_level = st.selectbox(
            "具体专业/水平", 
            ["大学英语（非英语专业大一）", "大学英语（非英语专业大二）", "英语专业低年级", "英语专业高年级/学术英语"]
        )
        
    default_duration = 45 if school_stage != "大学" else 90
    class_duration = st.number_input(
        "3. 课堂总时长（分钟）", 
        min_value=5, 
        max_value=180, 
        value=default_duration, 
        step=5
    )
    
    unit_topic = st.text_input("4. 单元主题 (Unit Topic)", value="My Family" if school_stage == "小学" else "Environmental Protection")
    
    teaching_method = st.selectbox(
        "5. 期望的课堂教学模式",
        [
            "💡 启发互动式课堂 (导入-讲解-互动提问-随堂反馈) [BOPPPS模型]",
            "🎯 任务学用一体式课堂 (通过实际应用产出任务来推动学习) [POA模型]",
            "🤝 情境任务探究式课堂 (设置小组合作和情境，在做中学会) [TBLT模型]",
            "📖 经典精讲多练式课堂 (传统的老师讲解、学生练习、最后输出) [3P模式]"
        ]
    )
    
    uploaded_file = st.file_uploader(
        "6. 上传课本课文/PPT/插图（支持 PDF、Word、PPTX、TXT、PNG、JPG）",
        type=["pdf", "docx", "pptx", "txt", "png", "jpg", "jpeg"]
    )
    
    range_label = "内容解析范围设置"
    unit_name = "页/张/段/行"
    max_val = 9999
    file_ext = ""
    file_bytes = b""
    
    use_vision_pdf = False # 默认不开启视觉扫描
    
    if uploaded_file is not None:
        file_ext = uploaded_file.name.split('.')[-1].lower()
        file_bytes = uploaded_file.getvalue()
        
        if file_ext == 'pdf':
            range_label = "📄 PDF 页码解析范围"
            unit_name = "页"
            try:
                pdf_reader = pypdf.PdfReader(BytesIO(file_bytes))
                max_val = len(pdf_reader.pages)
            except: pass
            
            # --- 常驻指南折叠卡与通俗易懂的勾选提示 ---
            st.write("---")
            with st.expander("❓ 如何选择是否开启“深度视觉扫描”？", expanded=True):
                st.markdown("""
                请观察您上传的 PDF 文件，做一次极简单的**文字复制测试**：
                1. ❌ **不勾选【默认关闭】（推荐 ⚡ 极速解析）**：
                   *   **判断方法**：在您平时的 PDF 阅读器中，页面上的英文或中文文字**可以用鼠标选中和复制**。
                   *   **适用场景**：电子版原装教材、文字版教案、纯文字课本。
                2. 👁️ **勾选【深度视觉扫描】（🔮 逐页读图解析）**：
                   *   **判断方法**：页面文字**无法用鼠标选中**（相当于一张张整页的照片，是用手机拍照或彩色扫描合成的 PDF）。
                   *   **特殊需求**：课本里有**大量彩色漫画、插图或排版图表**，您希望 AI 能够“看懂这些插画”并据此为您设计游戏和课堂提问。
                """)
            
            use_vision_pdf = st.checkbox("🔍 开启 PDF 深度视觉扫描 (读图解析)", value=False)
            if use_vision_pdf:
                # 替换掉了波浪号 ~ 避免被 markdown 语法错误识别成删除线
                st.warning("⚠️ 提醒：深度视觉扫描会逐页调用大模型进行解析，生成速度会慢一些（每页约 3 至 5 秒）。建议配合下方的“自定义范围”限制每次只扫描要讲的 2 至 3 页。")
            st.write("---")
            
        elif file_ext == 'pptx':
            range_label = "📊 PPT 幻灯片解析范围"
            unit_name = "张幻灯片"
            try:
                prs = Presentation(BytesIO(file_bytes))
                max_val = len(prs.slides)
            except: pass
        elif file_ext == 'docx':
            range_label = "📝 Word 自然段落解析范围"
            unit_name = "个段落"
            try:
                doc = docx.Document(BytesIO(file_bytes))
                max_val = len(doc.paragraphs)
            except: pass
        elif file_ext == 'txt':
            range_label = "🔤 TXT 文本行数解析范围"
            unit_name = "行"
            try:
                lines = file_bytes.decode("utf-8").split('\n')
                max_val = len(lines)
            except: pass
        elif file_ext in ['png', 'jpg', 'jpeg']:
            range_label = "🖼️ 图片解析范围"
            unit_name = "张图片"
            max_val = 1

    # 渲染范围控制器
    st.write(f"⚙️ **{range_label}**（当前文件共有 {max_val} {unit_name}）")
    range_type = st.radio("选择解析范围", ["全部", "自定义范围"], horizontal=True, key="range_type")
    
    start_num = 1
    end_num = max_val
    if range_type == "自定义范围" and file_ext not in ['png', 'jpg', 'jpeg']:
        p_col1, p_col2 = st.columns(2)
        with p_col1:
            start_num = st.number_input(f"起始{unit_name}", min_value=1, max_value=max_val, value=1, step=1)
        with p_col2:
            end_num = st.number_input(f"结束{unit_name}", min_value=start_num, max_value=max_val, value=min(start_num + 5, max_val), step=1)
    
    final_text = ""
    if uploaded_file is not None:
        # 非缓存区域安全的气泡提示
        if file_ext in ['png', 'jpg', 'jpeg'] or (file_ext == 'pdf' and use_vision_pdf):
            st.toast("🔮 通义千文正在解析视觉内容...", icon="👁️")
            
        final_text = extract_text_cached(
            uploaded_file.name, 
            file_bytes, 
            file_ext, 
            range_type, 
            start_num, 
            end_num,
            use_vision_pdf
        )
        if final_text and not final_text.startswith("ERROR"):
            with st.expander(f"👁️ 查看已成功解析的文字与画面描述 (共 {len(final_text)} 字)"):
                st.text_area("解析内容如下：", value=final_text[:5000] + ("\n...（已省略超出部分）" if len(final_text) > 5000 else ""), height=150, disabled=True)
        elif final_text.startswith("ERROR"):
            st.error(final_text)
            final_text = ""

    st.write("---")
    manual_text = st.text_area(
        "或者：手动输入/粘贴课文摘要或核心词汇 (选填)",
        placeholder="若未上传文件，可在此手动输入...",
        height=80
    )
    
    if not final_text:
        final_text = manual_text
        
    include_sizheng = st.checkbox("在教案中融入“思政育人/情感态度”模块", value=True)
    generate_btn = st.button("🚀 开始全新定制教案与大纲", type="primary", use_container_width=True)

with col2:
    st.subheader("📝 智能生成结果展示")
    
    # 5.1 【全生成】主链路运行
    if generate_btn:
        if not unit_topic:
            st.warning("请输入单元主题！")
        else:
            try:
                if class_duration <= 25:
                    ppt_pages_guideline = "5 - 8 张"
                    plan_depth = "精炼紧凑型教案。聚焦核心语言点，重点突出。"
                elif class_duration <= 50:
                    ppt_pages_guideline = "10 - 15 张"
                    plan_depth = "标准课时教学设计。包含完整的导入、呈现、讲授、课堂互动和作业布置，各环节时间分配精细。"
                else:
                    ppt_pages_guideline = "18 - 25 张"
                    plan_depth = "高容量、高深度教案。设计 2-3 个学生高阶产出任务，步骤丰富。"

                # 极其细致的多模态、高保真、无省略的系统提示词（升级：高对比度深色背景指令 + 真实 Unsplash 配图库）
                system_prompt = (
                    "你是一位极其严谨的高校与中小学英语教学研究专家，同时也是一位前沿幻灯片视觉设计师。\n"
                    "你需要完成一份详细教案、一套高排版美感且面向下游生成器无缝解析的【高保真 PPT 视觉大纲】、以及一组随堂测验题。\n\n"
                    "【⚠️ 核心对比度控制（解决白字在白底上看不清的问题）】\n"
                    "因为用户的下游 PPT 生成器默认使用您的全局色号，为了彻底解决文字与背景色温混淆的问题，请遵守以下设计规范：\n"
                    "- 所有幻灯片背景色统一设定为深海军蓝 (#1A365D)，以此在视觉上完美托出白色 (#FFFFFF) 或鲜亮黄 (#F1C40F) 的大标题，确保绝对对比度与易读性！\n"
                    "- 幻灯片内部的内容卡片容器（Containers）统一设定为纯白色 (#FFFFFF) 或浅亮色 (#F8F9FA)，卡片内部的正文文字设定为深灰色 (#333333) 或深海蓝 (#1A365D)，确保白底黑字，对比度极高！\n"
                    "绝对禁止在浅色背景上输出白色文字！\n\n"
                    "【⚠️ 真实插画配图规范（解决 PPT 没配图的问题）】\n"
                    "在每一张 Slide 卡片的 'Visual & Image Cue' 字段中，你必须直接输出以下列表中与当前页面主题最契合的高清、免版权 Unsplash 真实图片网络链接（使用标准 Markdown 图片格式 `![Description](URL)`），以便用户的 PPT 生成器可以直接提取并下载显示：\n"
                    "- 职场建筑/封面配图：`![Office](https://images.unsplash.com/photo-1486406146926-c627a92ad1ab?auto=format&fit=crop&w=600&q=80)`\n"
                    "- 商务精英/正式西装：`![Formal Suit](https://images.unsplash.com/photo-1560250097-0b93528c311a?auto=format&fit=crop&w=600&q=80)`\n"
                    "- 休闲装/硅谷程序员：`![Casual Wear](https://images.unsplash.com/photo-1519085360753-af0119f7cbe7?auto=format&fit=crop&w=600&q=80)`\n"
                    "- 职业规划/团队合作：`![Career](https://images.unsplash.com/photo-1522071820081-009f0129c71c?auto=format&fit=crop&w=600&q=80)`\n"
                    "- 词汇学习/专业讨论：`![Learning](https://images.unsplash.com/photo-1516321318423-f06f85e504b3?auto=format&fit=crop&w=600&q=80)`\n"
                    "- 结构分析/逻辑框架：`![Structure](https://images.unsplash.com/photo-1557804506-669a67965ba0?auto=format&fit=crop&w=600&q=80)`\n"
                    "- 随堂练习/笔头测试：`![Practice](https://images.unsplash.com/photo-1488190211105-8b0e65b80b4e?auto=format&fit=crop&w=600&q=80)`\n"
                    "- 迷你辩论/思想碰撞：`![Debate](https://images.unsplash.com/photo-1517245386807-bb43f82c33c4?auto=format&fit=crop&w=600&q=80)`\n"
                    "- 思政总结/职业精神：`![Professionalism](https://images.unsplash.com/photo-1454165804606-c3d57bc86b40?auto=format&fit=crop&w=600&q=80)`\n"
                    "- 课后作业/行动计划：`![Homework](https://images.unsplash.com/photo-1506784983877-45594efa4cbe?auto=format&fit=crop&w=600&q=80)`\n\n"
                    "【⚠️ 严格的无省略规范（No Ellipsis Guarantee）】\n"
                    "- 严禁使用任何“以此类推”、“此处省略”或占位符，每一张 Slide 的 Points Content 必须输出完整、详细、无缺失、可以直接朗读的演示文本！\n\n"
                    "请务必严格使用以下带有中括号的标记来分割你输出的内容（不可拼错）：\n"
                    "[LESSON_PLAN_START]\n"
                    "此处撰写详细教案（符合选定学段）。教学步骤中的时间累计必须正好等于总时长。如果使用表格展示时间分配和环节，请务必保证表格结构规整。\n"
                    "[LESSON_PLAN_END]\n\n"
                    "[OUTLINE_START]\n"
                    "此处撰写配套的【PPT 视觉设计大纲】。请严格按照以下结构化 Markdown 层次进行输出：\n\n"
                    "### 1. PPT Global Specifications (PPT全局视觉设计规范)\n"
                    "- **Design Style**: [设计风格风格描述，如：Modern Professional Grid / Minimalist Academic]\n"
                    "- **Background Color**: [整套 PPT 背景色 Hex 编码，强制设为 #1A365D]\n"
                    "- **Primary Theme Color**: [主色调 Hex 编码，推荐 #FFFFFF 或 #F1C40F]\n"
                    "- **Secondary Accent Color**: [辅助强调色 Hex 编码，推荐 #E67E22]\n"
                    "- **Default Text Color**: [正文字体颜色 Hex 编码，推荐 #333333 或 #1A365D]\n"
                    "- **Font Family**: [标题和正文字体，如：Arial, Microsoft YaHei]\n\n"
                    "### 2. Slide Cards (单页幻灯片内容与高精度版式方案)\n"
                    "请根据课堂时长，严格将幻灯片卡片总页数控制在 {ppt_pages_guideline} 页左右。\n"
                    "--- (使用三个减号作为每张 Slide 的开始分界线)\n"
                    "#### Slide 1: [单页幻灯片大标题]\n"
                    "- **Layout Recommended**: [极其具体的单页空间排版与栅格布局说明。例如：左右非对称分栏。左侧 55% 宽度为深蓝色块背景，内置白色大号加粗主标题与小号副标题；右侧 45% 为卡片式圆角容器 (#FFFFFF，圆角 12px)，内置核心矢量配图。]\n"
                    "- **Colors & Typography**: [此页具体的精细配色与字号方案。格式必须为：前景色-背景色成对标注，例如：幻灯片整体背景为 #1A365D。大标题使用 #FFFFFF (28pt/Bold)。卡片背景使用 #FFFFFF。卡片正文使用 #333333 (12pt)。重点强调词使用 #E67E22 (加粗)。确保前景色与背景色有极高对比度。]\n"
                    "- **Visual & Image Cue**: [必须从上方的真实图片数据库中，挑选最合适的一个 Markdown 图片格式链接直接输出。例如：`![Office](https://images.unsplash.com/photo-1486406146926-c627a92ad1ab?auto=format&fit=crop&w=600&q=80)`]\n"
                    "- **Points Content**:\n"
                    "  - [要点 1：必须是完整、无省略的演示文本，直接呈现在 PPT 上]\n"
                    "  - [要点 2：必须是完整、无省略的演示文本，直接呈现在 PPT 上]\n"
                    "  - [要点 3：必须是完整、无省略的演示文本，直接呈现在 PPT 上]\n"
                    "- **Teacher Notes**: [教师在这个幻灯片页面讲解时的完整手稿讲义、提问、以及过渡引导句]\n\n"
                    "--- (以此类推，继续写完后续所有幻灯片内容)\n"
                    "[OUTLINE_END]\n\n"
                    "[QUIZ_START]\n"
                    "此处撰写适合该学段和课本内容的随堂测验与课堂互动方案：\n"
                    "1. 精心设计 3 道极有针对性的词汇或语法单选题（附标准答案与简要解析）。\n"
                    "2. 设计 1 道与本课文本相关的阅读理解题或填空题。\n"
                    "3. 专为该学段设计 1 个课堂游戏或趣味互动竞技活动（具体操作步骤和奖励机制说明）。\n"
                    "[QUIZ_END]\n"
                )
                
                user_prompt = f"""
                请为我生成英语教案、大纲与随堂小练：
                - 课程/科目名称: {course_name}
                - 授课学段: {school_stage}（具体为：{grade_level}）
                - 课堂总时长: {class_duration} 分钟
                - 教学深度要求: {plan_depth}
                - PPT 课件大纲篇幅要求: 约 {ppt_pages_guideline} 页左右
                - 单元/课文主题: {unit_topic}
                - 期望的课堂模式: {teaching_method}
                - 思政育人/情感态度融入: {"需要深入融入该学段专属的思政育人方向" if include_sizheng else "暂不融合"}
                
                【课本实际内容/课文选段与画面解析】：
                {final_text if final_text else "（未提供具体课本，请围绕主题进行该学段的通用高水准教学设计）"}
                """
                
                st.info(f"⚡ AI 正在实时为您撰写整体教学包方案（含高精度 PPT 设计元数据）...")
                stream_placeholder = st.empty()
                raw_result = ""
                
                response = client_deepseek.chat.completions.create(
                    model="deepseek-chat",
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt}
                    ],
                    stream=True
                )
                
                for chunk in response:
                    if chunk.choices[0].delta.content is not None:
                        raw_result += chunk.choices[0].delta.content
                        stream_placeholder.markdown(raw_result)
                
                # 数据分割
                lesson_plan_data = ""
                outline_data = ""
                quiz_data = ""
                
                if "[LESSON_PLAN_START]" in raw_result and "[LESSON_PLAN_END]" in raw_result:
                    lesson_plan_data = raw_result.split("[LESSON_PLAN_START]")[1].split("[LESSON_PLAN_END]")[0].strip()
                else:
                    lesson_plan_data = raw_result
                    
                if "[OUTLINE_START]" in raw_result and "[OUTLINE_END]" in raw_result:
                    outline_data = raw_result.split("[OUTLINE_START]")[1].split("[OUTLINE_END]")[0].strip()
                else:
                    outline_data = "AI 未能成功生成独立大纲，请参考教案内容。"
                    
                if "[QUIZ_START]" in raw_result and "[QUIZ_END]" in raw_result:
                    quiz_data = raw_result.split("[QUIZ_START]")[1].split("[QUIZ_END]")[0].strip()
                else:
                    quiz_data = "AI 未能成功生成随堂测试，请再次尝试或微调重新生成。"
                
                stream_placeholder.empty()
                st.session_state['generated_plan'] = lesson_plan_data
                st.session_state['generated_outline'] = outline_data
                st.session_state['generated_quiz'] = quiz_data
                st.success("🎉 教学包（教案 + 视觉大纲 + 随堂小练）生成完毕！")
                
            except Exception as e:
                st.error(f"处理出错: {e}")

    # 5.2 【微调】二次对话微调运行区域
    st.write("---")
    st.write("💬 **“智能副驾驶”二次局部微调（Tweak Engine）**")
    tweak_text = st.text_input(
        "在此输入修改指令，对上方已生成的内容进行针对性微改：",
        placeholder="例如：'把PPT大纲的主色色号调整为薄荷绿'、'让第3张Slide的排版变为居中三列卡片'",
        key="tweak_input"
    )
    tweak_btn = st.button("🔧 提交修改指令", use_container_width=True)

    if tweak_btn and tweak_text:
        if not st.session_state['generated_plan']:
            st.warning("请先在上方成功生成一次教案，然后才能提出修改指令！")
        else:
            try:
                tweak_system_prompt = (
                    "你是一个资深英语教学和课件大纲设计专家。现在，用户对刚刚生成的教学包（包含：详细教案、授课大纲、随堂随练）提出了修改意见。\n"
                    "请你根据用户的修改建议，对原有的三部分内容进行精准微调或重写。\n\n"
                    "为了保证系统能正确解析和分栏显示，请您**必须完整输出微调后的全部三部分内容**，并依然严格套用中括号标记进行分隔：\n"
                    "[LESSON_PLAN_START]\n微调后的完整教案...\n[LESSON_PLAN_END]\n\n"
                    "[OUTLINE_START]\n微调后的完整授课大纲（带设计元数据、配色色号、字体和详细插画排版描述）...\n[OUTLINE_END]\n\n"
                    "[QUIZ_START]\n微调后的完整随堂随练与活动设计...\n[QUIZ_END]\n"
                )
                
                tweak_user_prompt = f"""
                原教学包数据如下：
                ====================
                [原详细教案]
                {st.session_state['generated_plan']}
                
                [原授课大纲]
                {st.session_state['generated_outline']}
                
                [原随堂测试与活动]
                {st.session_state['generated_quiz']}
                ====================
                
                用户的修改建议：\"{tweak_text}\"
                
                请严格按照修改建议，重新流式输出微改后的完整教学设计方案。
                """
                
                st.info("⚡ 副驾驶正在按您的要求调整教案方案，请预览微调重写过程...")
                tweak_placeholder = st.empty()
                tweak_raw_result = ""
                
                response_tweak = client_deepseek.chat.completions.create(
                    model="deepseek-chat",
                    messages=[
                        {"role": "system", "content": tweak_system_prompt},
                        {"role": "user", "content": tweak_user_prompt}
                    ],
                    stream=True
                )
                
                for chunk in response_tweak:
                    if chunk.choices[0].delta.content is not None:
                        tweak_raw_result += chunk.choices[0].delta.content
                        tweak_placeholder.markdown(tweak_raw_result)
                
                # 微调数据切分
                tweak_plan_data = ""
                tweak_outline_data = ""
                tweak_quiz_data = ""
                
                if "[LESSON_PLAN_START]" in tweak_raw_result and "[LESSON_PLAN_END]" in tweak_raw_result:
                    tweak_plan_data = tweak_raw_result.split("[LESSON_PLAN_START]")[1].split("[LESSON_PLAN_END]")[0].strip()
                else:
                    tweak_plan_data = tweak_raw_result
                    
                if "[OUTLINE_START]" in tweak_raw_result and "[OUTLINE_END]" in tweak_raw_result:
                    tweak_outline_data = tweak_raw_result.split("[OUTLINE_START]")[1].split("[OUTLINE_END]")[0].strip()
                else:
                    tweak_outline_data = st.session_state['generated_outline']
                    
                if "[QUIZ_START]" in tweak_raw_result and "[QUIZ_END]" in tweak_raw_result:
                    tweak_quiz_data = tweak_raw_result.split("[QUIZ_START]")[1].split("[QUIZ_END]")[0].strip()
                else:
                    tweak_quiz_data = st.session_state['generated_quiz']
                
                tweak_placeholder.empty()
                st.session_state['generated_plan'] = tweak_plan_data
                st.session_state['generated_outline'] = tweak_outline_data
                st.session_state['generated_quiz'] = tweak_quiz_data
                st.success("🎉 修改意见已采纳并更新完成！")
                
            except Exception as e:
                st.error(f"微调失败: {e}")

    # 5.3 【渲染分栏】始终根据最新的 Session State 保持分栏渲染
    tab1, tab2, tab3 = st.tabs([
        "📄 详细教案 (Lesson Plan)", 
        "📝 PPT 授课大纲 (Slide Outline)",
        "✏️ 随堂随练 (Quiz & Activities)"
    ])
    
    if st.session_state['generated_plan'] or st.session_state['generated_outline'] or st.session_state['generated_quiz']:
        with tab1:
            st.markdown(st.session_state['generated_plan'])
        with tab2:
            st.markdown(st.session_state['generated_outline'])
        with tab3:
            st.markdown(st.session_state['generated_quiz'])
            
        st.write("---")
        
        # 将教案、大纲、小练整合成一个 Word 文件下载
        docx_file = convert_markdown_to_docx(
            st.session_state['generated_plan'], 
            st.session_state['generated_outline'], 
            st.session_state['generated_quiz']
        )
        
        st.download_button(
            label="📥 下载精美排版 Word 版（教案+大纲+随堂小练）(.docx)",
            data=docx_file,
            file_name=f"{course_name}_{unit_topic}_完整教学包.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            use_container_width=True
        )
    else:
        st.info("请在左侧配置参数并上传课本，点击“开始定制教案与大纲”按钮。")
